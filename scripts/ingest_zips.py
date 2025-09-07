#!/usr/bin/env python3
# saves text chunks to ~/tullman/data/content/content.jsonl
# saves PPT images to ~/tullman/data/media and logs them in media_manifest.jsonl
import argparse, json, re, zipfile, uuid, hashlib
from pathlib import Path
from tqdm import tqdm
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import chardet

def sha(s): return hashlib.sha256(s.encode("utf-8","ignore")).hexdigest()
def rd(p): return open(p,"rb").read()
def guess_text(p):
    b = rd(p)
    try: return b.decode("utf-8")
    except: return b.decode(chardet.detect(b).get("encoding") or "latin-1", errors="ignore")
def from_pdf(p):
    out=[]; r=PdfReader(str(p))
    for pg in r.pages: out.append(pg.extract_text() or "")
    return "\n".join(out)
def from_docx(p):
    d=Document(str(p)); return "\n".join(x.text for x in d.paragraphs)
def from_html(p):
    s=BeautifulSoup(guess_text(p), "lxml")
    for t in s(["script","style"]): t.extract()
    return s.get_text(separator="\n")
def chunk(s, maxc=1200, ov=120):
    s=re.sub(r"\n{3,}","\n\n",s).strip()
    out=[]; i=0
    while i < len(s):
        j=min(i+maxc,len(s))
        cut=s.rfind("\n\n",i,j); 
        if cut==-1: cut=s.rfind(". ",i,j)
        if cut==-1: cut=j
        out.append(s[i:cut].strip()); i=max(cut-ov,cut)
    return [x for x in out if x]
def write_jsonl(path, obj):
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path,"a",encoding="utf-8") as f: f.write(json.dumps(obj,ensure_ascii=False)+"\n")
def load_seen(path):
    seen=set()
    if path.exists():
        for line in open(path,"r",encoding="utf-8"):
            try:
                h=json.loads(line).get("hash"); 
                if h: seen.add(h)
            except: pass
    return seen
def harvest_ppt_media(pptx_path, media_dir, manifest):
    saved=0
    with zipfile.ZipFile(str(pptx_path),"r") as z:
        for name in z.namelist():
            if name.startswith("ppt/media/") and not name.endswith("/"):
                data=z.read(name)
                out = media_dir / pptx_path.stem / Path(name).name
                out.parent.mkdir(parents=True, exist_ok=True)
                with open(out,"wb") as f: f.write(data)
                write_jsonl(manifest, {"file":str(out.relative_to(media_dir)), "title":f"{pptx_path.stem} asset", "origin":str(pptx_path)})
                saved+=1
    return saved
def slide_texts(p):
    prs=Presentation(str(p)); out=[]
    for i,slide in enumerate(prs.slides,1):
        bits=[]
        for sh in slide.shapes:
            try:
                if hasattr(sh,"text") and sh.text: bits.append(sh.text)
                elif getattr(sh,"has_text_frame",False) and sh.text_frame: bits.append(sh.text_frame.text)
            except: pass
        t="\n".join([b for b in bits if b and b.strip()])
        out.append((i,t.strip()))
    return out

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("--zips", nargs="+", required=True)
    base = Path("~/tullman").expanduser()
    ap.add_argument("--raw", default=str(base/"data/raw"))
    ap.add_argument("--content", default=str(base/"data/content/content.jsonl"))
    ap.add_argument("--media", default=str(base/"data/media"))
    ap.add_argument("--manifest", default=str(base/"data/media_manifest.jsonl"))
    a=ap.parse_args()

    raw=Path(a.raw); raw.mkdir(parents=True, exist_ok=True)
    media=Path(a.media); media.mkdir(parents=True, exist_ok=True)
    content=Path(a.content); manifest=Path(a.manifest)
    seen=load_seen(content)

    # unzip
    roots=[]
    for z in a.zips:
        zp=Path(z).expanduser()
        if not zp.exists(): print(f"[skip] {zp}"); continue
        out=raw/zp.stem; out.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(str(zp),"r") as zf: zf.extractall(str(out))
        roots.append(out)

    newc=0; media_ct=0
    for root in roots:
        files=[p for p in root.rglob("*") if p.is_file()]
        for p in tqdm(files, desc=f"scan {root.name}"):
            ext=p.suffix.lower()
            if ext==".pdf":
                txt=from_pdf(p); parts=chunk(txt); st="pdf"
            elif ext==".docx":
                txt=from_docx(p); parts=chunk(txt); st="docx"
            elif ext in {".txt",".md"}:
                txt=guess_text(p); parts=chunk(txt); st=ext.strip(".")
            elif ext in {".html",".htm"}:
                txt=from_html(p); parts=chunk(txt); st="html"
            elif ext==".pptx":
                media_ct+=harvest_ppt_media(p, media, manifest)
                for idx, s in slide_texts(p):
                    if not s: continue
                    h=sha(f"{p}::slide::{idx}::{s[:400]}")
                    if h in seen: continue
                    seen.add(h)
                    write_jsonl(content, {
                        "id": str(uuid.uuid4()),
                        "title": f"{p.name} slide {idx}",
                        "source_path": str(p),
                        "source_name": p.name,
                        "source_type": "pptx",
                        "part": f"slide_{idx}",
                        "text": s,
                        "tags": ["tullman_ai","slides"],
                        "hash": h
                    }); newc+=1
                continue
            else:
                continue
            for i, c in enumerate(parts,1):
                h=sha(f"{p}::chunk::{i}::{c[:400]}")
                if h in seen: continue
                seen.add(h)
                write_jsonl(content, {
                    "id": str(uuid.uuid4()),
                    "title": f"{p.name} chunk {i}",
                    "source_path": str(p),
                    "source_name": p.name,
                    "source_type": st,
                    "part": f"chunk_{i}",
                    "text": c,
                    "tags": ["tullman_ai"],
                    "hash": h
                }); newc+=1
    print(f"[done] new content: {newc} | media saved: {media_ct}")
    print(f"[paths] {content}\n[paths] {manifest}\n[paths] {media}")

if __name__=="__main__":
    main()
