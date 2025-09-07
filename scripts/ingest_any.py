#!/usr/bin/env python3
"""
Tolerant ingester with streaming progress:
- Accepts inputs that are .zip files, directories, or single files (pdf, docx, txt, md, html, htm, pptx).
- Skips oversized files with --max-file-mb.
- Streams progress every --log-every files.
- Smarter chunking (paragraph/word based) avoids slow rfind scanning on huge strings.
Writes:
  ~/tullman/data/content/content.jsonl
  ~/tullman/data/media_manifest.jsonl
"""
import argparse, json, uuid, hashlib, zipfile, re, sys
from pathlib import Path
from typing import Iterable, List, Set

# Lightweight parsers
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import chardet

BASE = Path("~/tullman").expanduser()
RAW_DIR = BASE / "data" / "raw"
CONTENT_JSONL = BASE / "data" / "content" / "content.jsonl"
MEDIA_DIR = BASE / "data" / "media"
MEDIA_MANIFEST = BASE / "data" / "media_manifest.jsonl"

SUPPORTED_TEXT_EXT = {".pdf", ".docx", ".txt", ".md", ".html", ".htm", ".pptx"}

def sha(s:str)->str:
    return hashlib.sha256(s.encode("utf-8","ignore")).hexdigest()

def write_jsonl(path:Path, obj:dict):
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(obj, ensure_ascii=False) + "\n")

def load_seen(path:Path)->Set[str]:
    seen=set()
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                try:
                    j=json.loads(line); h=j.get("hash")
                    if h: seen.add(h)
                except: pass
    return seen

def read_bytes(p:Path)->bytes:
    with open(p, "rb") as f: return f.read()

def guess_text(p:Path)->str:
    data=read_bytes(p)
    try: return data.decode("utf-8")
    except UnicodeDecodeError:
        enc = chardet.detect(data).get("encoding") or "latin-1"
        return data.decode(enc, errors="ignore")

def text_from_pdf(p:Path)->str:
    out=[]
    try:
        r=PdfReader(str(p))
        for pg in r.pages: out.append(pg.extract_text() or "")
    except Exception as e:
        out.append(f"[pdf read error: {e}]")
    return "\n".join(out)

def text_from_docx(p:Path)->str:
    try:
        d=Document(str(p))
        return "\n".join(para.text for para in d.paragraphs)
    except Exception as e:
        return f"[docx read error: {e}]"

def text_from_html(p:Path)->str:
    try:
        soup=BeautifulSoup(guess_text(p), "lxml")
        for t in soup(["script","style"]): t.extract()
        return soup.get_text(separator="\n")
    except Exception as e:
        return f"[html parse error: {e}]"

def chunk_smart(text:str, max_chars:int=1200, overlap:int=120)->List[str]:
    """
    Fast, paragraph-first chunker. No repeated rfind scans.
    1) Split on blank lines into paragraphs.
    2) Pack paragraphs into chunks up to max_chars.
    3) If there were no paragraphs, fall back to word-based packing.
    """
    if not text: return []
    # normalize whitespace a bit
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text: return []

    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks=[]
    if paras:
        cur=[]
        cur_len=0
        for p in paras:
            # +2 for the double newline we removed
            add_len = (2 if cur else 0) + len(p)
            if cur_len + add_len <= max_chars or not cur:
                cur.append(p); cur_len += add_len
            else:
                chunks.append("\n\n".join(cur))
                # start next chunk with overlap from end of previous if helpful
                if overlap and chunks[-1]:
                    tail = chunks[-1][-overlap:]
                    cur=[tail, p]; cur_len=len(tail)+2+len(p)
                else:
                    cur=[p]; cur_len=len(p)
        if cur: chunks.append("\n\n".join(cur))
    else:
        # word packing
        words = text.split()
        cur=[]; cur_len=0
        for w in words:
            add_len = (1 if cur else 0) + len(w)
            if cur_len + add_len <= max_chars or not cur:
                cur.append(w); cur_len += add_len
            else:
                chunks.append(" ".join(cur))
                if overlap and chunks[-1]:
                    tail = chunks[-1][-overlap:]
                    cur=[tail, w]; cur_len=len(tail)+1+len(w)
                else:
                    cur=[w]; cur_len=len(w)
        if cur: chunks.append(" ".join(cur))
    return [c for c in chunks if c.strip()]

def harvest_pptx_media(pptx_path:Path):
    """Save embedded media to MEDIA_DIR/<deck>/ and append manifest rows."""
    saved=0
    try:
        with zipfile.ZipFile(str(pptx_path), "r") as z:
            for name in z.namelist():
                if name.startswith("ppt/media/") and not name.endswith("/"):
                    data = z.read(name)
                    out = MEDIA_DIR / pptx_path.stem / Path(name).name
                    out.parent.mkdir(parents=True, exist_ok=True)
                    with open(out, "wb") as f: f.write(data)
                    write_jsonl(MEDIA_MANIFEST, {
                        "file": str(out.relative_to(MEDIA_DIR)),
                        "title": f"{pptx_path.stem} asset",
                        "origin": str(pptx_path)
                    })
                    saved += 1
    except Exception:
        pass
    return saved

def slide_texts(p:Path):
    out=[]
    try:
        prs=Presentation(str(p))
        for i, slide in enumerate(prs.slides, start=1):
            bits=[]
            for sh in slide.shapes:
                try:
                    if hasattr(sh, "text") and sh.text:
                        bits.append(sh.text)
                    elif getattr(sh, "has_text_frame", False) and sh.text_frame:
                        bits.append(sh.text_frame.text)
                except Exception:
                    pass
            txt="\n".join([b for b in bits if b and b.strip()]).strip()
            out.append((i, txt))
    except Exception as e:
        out.append((0, f"[pptx parse error: {e}]"))
    return out

def scan_dir(d:Path)->Iterable[Path]:
    for p in d.rglob("*"):
        if p.is_file(): yield p

def process_file(p:Path, seen:Set[str], counters:dict, args):
    ext = p.suffix.lower()
    base_title = p.name

    # size guard
    try:
        size_mb = p.stat().st_size / (1024*1024)
        if size_mb > args.max_file_mb:
            print(f"[skip-oversize] {p} ({size_mb:.1f} MB > {args.max_file_mb} MB)", flush=True)
            return
    except Exception:
        pass

    # type filter (optional)
    if args.types:
        if ext.lstrip(".") not in args.types and ext not in {f".{t}" for t in args.types}:
            return

    if ext == ".pdf":
        text = text_from_pdf(p); parts = chunk_smart(text, args.max_chars, args.overlap); src_type="pdf"
    elif ext == ".docx":
        text = text_from_docx(p); parts = chunk_smart(text, args.max_chars, args.overlap); src_type="docx"
    elif ext in {".txt",".md"}:
        text = guess_text(p); parts = chunk_smart(text, args.max_chars, args.overlap); src_type=ext.strip(".")
    elif ext in {".html",".htm"}:
        text = text_from_html(p); parts = chunk_smart(text, args.max_chars, args.overlap); src_type="html"
    elif ext == ".pptx":
        counters["media"] += harvest_pptx_media(p)
        for idx, s in slide_texts(p):
            if not s: continue
            h = sha(f"{p}::slide::{idx}::{s[:400]}")
            if h in seen: continue
            seen.add(h)
            write_jsonl(CONTENT_JSONL, {
                "id": str(uuid.uuid4()),
                "title": f"{base_title} slide {idx}",
                "source_path": str(p),
                "source_name": base_title,
                "source_type": "pptx",
                "part": f"slide_{idx}",
                "text": s,
                "tags": ["tullman_ai","slides"],
                "hash": h
            })
            counters["content"] += 1
        return
    else:
        return  # silently ignore unknown types

    for i, chunk in enumerate(parts, start=1):
        if not chunk.strip(): continue
        h = sha(f"{p}::chunk::{i}::{chunk[:400]}")
        if h in seen: continue
        seen.add(h)
        write_jsonl(CONTENT_JSONL, {
            "id": str(uuid.uuid4()),
            "title": f"{base_title} chunk {i}",
            "source_path": str(p),
            "source_name": base_title,
            "source_type": src_type,
            "part": f"chunk_{i}",
            "text": chunk,
            "tags": ["tullman_ai"],
            "hash": h
        })
        counters["content"] += 1

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--inputs", nargs="+", required=True,
                    help="Paths to zip files, directories, or single files")
    ap.add_argument("--max-file-mb", type=float, default=80.0, help="Skip files larger than this")
    ap.add_argument("--log-every", type=int, default=25, help="Print a progress line every N files")
    ap.add_argument("--max-chars", type=int, default=1200, help="Max chars per chunk")
    ap.add_argument("--overlap", type=int, default=120, help="Overlap characters between chunks")
    ap.add_argument("--types", type=str, default="", help="Limit to types, e.g. 'pdf,docx'")
    args = ap.parse_args()
    args.types = [t.strip().lower() for t in args.types.split(",") if t.strip()] or None

    RAW_DIR.mkdir(parents=True, exist_ok=True)
    MEDIA_DIR.mkdir(parents=True, exist_ok=True)
    seen = load_seen(CONTENT_JSONL)
    counters = {"content": 0, "media": 0, "files": 0}

    to_scan = []

    for item in args.inputs:
        p = Path(item).expanduser()
        if not p.exists():
            print(f"[skip] not found: {p}", flush=True)
            continue

        if p.is_dir():
            to_scan.extend(list(scan_dir(p)))
            continue

        ext = p.suffix.lower()

        if ext == ".zip":
            try:
                out = RAW_DIR / p.stem
                out.mkdir(parents=True, exist_ok=True)
                with zipfile.ZipFile(str(p), "r") as zf:
                    zf.extractall(str(out))
                print(f"[unzipped] {p} -> {out}", flush=True)
                to_scan.extend(list(scan_dir(out)))
            except zipfile.BadZipFile:
                print(f"[warn] Not a valid zip: {p}. Skipping.", flush=True)
            continue

        if ext in SUPPORTED_TEXT_EXT:
            to_scan.append(p)
        else:
            print(f"[skip] unsupported file type: {p}", flush=True)

    total = len(to_scan)
    print(f"[scan] {total} files queued", flush=True)
    for f in to_scan:
        try:
            process_file(f, seen, counters, args)
        except KeyboardInterrupt:
            print("\n[abort] interrupted by user", flush=True)
            sys.exit(130)
        except Exception as e:
            print(f"[warn] failed on {f}: {e}", flush=True)
        counters["files"] += 1
        if counters["files"] % args.log_every == 0:
            print(f"[progress] {counters['files']}/{total} files • "
                  f"{counters['content']} content • {counters['media']} media", flush=True)

    print(f"[done] new content entries: {counters['content']}", flush=True)
    print(f"[done] media files harvested: {counters['media']}", flush=True)
    print(f"[paths] content: {CONTENT_JSONL}", flush=True)
    print(f"[paths] media manifest: {MEDIA_MANIFEST}", flush=True)
    print(f"[paths] media dir: {MEDIA_DIR}", flush=True)

if __name__ == "__main__":
    main()
